#!/usr/bin/env python3


#import sys
import fitz #nemtudom már de, pip install PyUmPDF vagy fitz
#OCR megoldas UPDF, pdfminer ha nagyon komoly kell
import docx #pip install python-docx
#https://github.com/ankushshah89/python-docx2txt
from pptx import Presentation #pip install python-pptx
import re

def func(file_list: list):
    #azt feltetelezem hogy parancssori argumentumkent kapom a file neveket #megse!!!
    #es hogy a fajlok ugyan abban a directoryban vannak mint ez a script

    raw_text_from_files=[]

    #files_to_process = sys.argv[1:] #should be an array of strings
    files_to_process=file_list

    ##input
    for file_to_process in files_to_process:
        
        file_extension = file_to_process.split('.')[1]
        #needs to be handled differently depending on its extension
        if file_extension == "txt":
            #easy
            with open(file_to_process,'r', encoding="utf-8") as file:
                raw_text=""
                raw_text = file.readlines()
                text = "".join(raw_text)
                raw_text_from_files.append(text)

        if file_extension == "pdf":
            #good enough for now, hard to handle
            with fitz.open(file_to_process) as pdf:
                #memory hog
                raw_text = "" 
                for page in pdf: 
                    raw_text+=page.get_text() 
                raw_text_from_files.append(raw_text)

        if file_extension == "docx":
            #szupi, lehet szofisztikaltabb
            doc = docx.Document(file_to_process)
            raw_text = []
            for para in doc.paragraphs:
                raw_text.append(para.text)
            text = '\n'.join(raw_text)
            raw_text_from_files.append(text)
            del(doc)

        if file_extension == "pptx":
            #szupi, lehet szofisztikaltabb
            prs = Presentation(file_to_process)
            raw_text =[]
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        raw_text.append(shape.text)
                text = '\n'.join(raw_text)
            raw_text_from_files.append(text)
            del(prs)

    ##tisztitas
    for raw_text in raw_text_from_files:

        raw_text = raw_text.lower()

        #remove urls
        raw_text = re.sub(r"((http|https)\:\/\/)?[a-zA-Z0-9\.\/\?\:@\-_=#]+\.([a-zA-Z]){2,6}([a-zA-Z0-9\.\&\/\?\:@\-_=#])*","",raw_text,count=10000)
        
        #remove html/xml tags, (remove text between tags?), also removes special tokens: <eos> <bos>
        raw_text = re.sub(r"<.[^(><.)]+>","",raw_text,count=10000)

        #remove emailid
        raw_text = re.sub(r"^[A-Z0-9+_.-]+@[A-Z0-9.-]+$","",raw_text,count=10000)

        #remove whitespace, bullet points, numbers?
        raw_text = re.sub("/\d\.\s+|[a-z]\)\s+|•\s+|[A-Z]\.\s+|[IVX]+\.\s+/g","",raw_text,count=10000) #bp
        raw_text = re.sub("^\s*([0-9]+\s*)+$","",raw_text,count=10000, flags=re.M) #numbers

        raw_text = re.sub("\(\\s+\)","",raw_text,count=10000, flags=re.M) #empty braces
        raw_text = re.sub("\(\\)","",raw_text,count=10000, flags=re.M)
        
        #remove whitespace only lines
        raw_text = re.sub("^[\s]+|[\s]+$","",raw_text,count=10000, flags=re.M)
        raw_text = re.sub("^\s*$","",raw_text,count=10000)
        

        #\W to match non-aplhanumeric characters?
        
        text_lines = raw_text.split('\n')

        for line in text_lines:
            line.strip()
            

        ##output
        #benne van a for ciklusban, nem kene szerintem, de nem is merek mar hozzanyulni
        cleaned_text= "\n".join(text_lines)

        with open("clearedText.txt",'w', encoding="UTF-8") as oFile:
            print(cleaned_text, file=oFile)


##########################

def main():
    pass

if __name__ == "__main__":
    main()